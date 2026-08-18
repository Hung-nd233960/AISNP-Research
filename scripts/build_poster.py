"""Fill in reports/TEMPLATE POSTER BME DESIGN 2025.2.pptx with this project's
actual content, writing reports/POSTER.pptx.

Strategy: edit the template in place rather than rebuild it, so every
decorative shape (colored panels, section-header ribbons, fonts) stays
exactly as designed. Two kinds of edits:

  1. Text placeholders ("Text", "Figure N: Text", "TOPIC: ", etc.) — the
     first run's font (Cambria, size, italic, color) is preserved; only the
     text itself changes. Extra lines reuse that same font.
  2. Picture placeholders — each lives alone inside its own single-child
     GROUP shape. The group is removed and replaced with a plain picture at
     the group's own bounding box, aspect-fit (never distorted/stretched).

The template's 6 picture slots are numbered Figure 1-6 in reading order
(methodology-left, methodology-right, research-purpose, testing-results,
then the two "prototype" slots, repurposed here as "Recommendation" since
this is a software/ML project, not a hardware prototype). That numbering is
local to the poster and intentionally does not match the report's Fig. 1-7.
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from PIL import Image

REPO = Path(__file__).resolve().parent.parent
FIG = REPO / "reports/figures"
TEMPLATE = REPO / 'reports/TEMPLATE POSTER BME DESIGN 2025.2.pptx'
OUT = REPO / "reports/POSTER.pptx"

AUTHOR = "Nguyễn Đức Hùng - 20233960"


def find(shapes, name):
    for s in shapes:
        if s.name == name:
            return s
        if s.shape_type == 6:
            r = find(s.shapes, name)
            if r is not None:
                return r
    return None


def set_text(shape, lines):
    """Overwrite a placeholder's text, preserving the first run's font,
    reusing that same font for any additional lines."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    base_run = p0.runs[0]
    font = base_run.font
    name, size, italic, bold = font.name, font.size, font.italic, font.bold
    color = font.color.rgb if font.color and font.color.type else None

    base_run.text = lines[0]
    for extra_run in list(p0.runs[1:]):
        extra_run._r.getparent().remove(extra_run._r)
    for extra_p in list(tf.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)

    for line in lines[1:]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = name
        if size:
            r.font.size = size
        r.font.italic = italic
        r.font.bold = bold
        if color:
            r.font.color.rgb = color


def replace_picture_group(slide, group_name, image_path):
    """Remove the single-child group `group_name` and add a plain,
    aspect-fit picture at the group's own bounding box."""
    group = find(slide.shapes, group_name)
    left, top, width, height = group.left, group.top, group.width, group.height

    with Image.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = width / height
    if aspect > box_aspect:
        w, h = width, int(width / aspect)
    else:
        h, w = height, int(height * aspect)
    new_left = left + (width - w) // 2
    new_top = top + (height - h) // 2

    group._element.getparent().remove(group._element)
    slide.shapes.add_picture(str(image_path), new_left, new_top, width=w, height=h)


prs = Presentation(TEMPLATE)
slide = prs.slides[0]

# ---------------------------------------------------------------------------
# Title / authorship
# ---------------------------------------------------------------------------
set_text(find(slide.shapes, "TextBox 26"),
         ["TOPIC: Selection of Ancestry-Informative SNP Panels for "
          "Within-East-Asian Population Discrimination"])
set_text(find(slide.shapes, "TextBox 92"), [f"Team member: {AUTHOR}", "Instructor: "])

# ---------------------------------------------------------------------------
# Body text
# ---------------------------------------------------------------------------
set_text(find(slide.shapes, "TextBox 19"), [  # ABSTRACT
    "A leak-free, automated pipeline selects compact SNP panels "
    "distinguishing Han, Japanese, and Southeast Asian ancestry. "
    "A 70-SNP panel reaches 92.9-94.25% accuracy, matching an "
    "expert-curated 34-SNP panel (94.6%)."
])
set_text(find(slide.shapes, "TextBox 18"), [  # INTRODUCTION
    "Han, Japanese, and Southeast Asian populations share recent common "
    "ancestry, so allele-frequency differences are small and diffuse, "
    "not concentrated in a few strong markers."
])
set_text(find(slide.shapes, "TextBox 34"), [  # RESEARCH PURPOSE & OBJECT
    "Build and validate a pipeline that selects and scores compact AISNP "
    "panels (N = 35, 50, 70) with zero leakage of test-sample information."
])
set_text(find(slide.shapes, "TextBox 75"), [  # TESTING AND RESULTS
    "Blind Selection accuracy rises from 65% (N=5) to 94% (N=100), "
    "validating the search before any panel is committed."
])
set_text(find(slide.shapes, "TextBox 91"), [  # RECOMMENDATION body
    "ElasticNet always; fst_stat as the default pool, trusting blind "
    "selection over hindsight; classifier pre-specified, not picked post hoc."
])
set_text(find(slide.shapes, "TextBox 63"), [  # FUTURE WORK
    "Validate on an independent East Asian cohort; widen the search "
    "beyond 54 candidate configurations."
])
set_text(find(slide.shapes, "TextBox 64"), [  # REFERENCES
    "Cai et al., Forensic Sci. Int. 357 (2024).",
    "1000 Genomes Project Consortium, Nature 526 (2015).",
])

# ---------------------------------------------------------------------------
# Section header rename: this template was built for a hardware prototype;
# repurposed here as the headline recommendation.
# ---------------------------------------------------------------------------
set_text(find(slide.shapes, "TextBox 84"), ["RECOMMENDATION"])

# ---------------------------------------------------------------------------
# Figures (poster-local numbering, reading order) + captions
# ---------------------------------------------------------------------------
replace_picture_group(slide, "Group 45", REPO / "Infographic.png")
set_text(find(slide.shapes, "TextBox 79"),
         ["Figure 1: Blind Selection vs. Best Fixed Configuration"])

replace_picture_group(slide, "Group 47", FIG / "fig2_panelfinder_detail.png")
set_text(find(slide.shapes, "TextBox 80"), ["Figure 2: Panel Finder detail"])

replace_picture_group(slide, "Group 35", FIG / "fig1_system_overview.png")
set_text(find(slide.shapes, "TextBox 37"), ["Figure 3: Overall system diagram"])

replace_picture_group(slide, "Group 76", FIG / "fig6_accuracy_vs_n.png")
set_text(find(slide.shapes, "TextBox 78"), ["Figure 4: Accuracy vs. panel size"])

replace_picture_group(slide, "Group 85", FIG / "fig7_confusion_matrices.png")
set_text(find(slide.shapes, "TextBox 87"), ["Figure 5: Confusion matrices"])

replace_picture_group(slide, "Group 88", FIG / "fig5_pool_divergence.png")
set_text(find(slide.shapes, "TextBox 90"), ["Figure 6: Pool divergence, hindsight vs. blind"])

prs.save(OUT)
print(f"Wrote {OUT}")
