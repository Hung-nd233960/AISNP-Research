"""Build reports/SLIDES.pptx — an editable PowerPoint deck distilled from
the finalized reports/REPORT.md. Mirrors reports/SLIDES.md.

Short text, diagram-heavy: every slide that has a figure in the report leans
on that figure instead of restating it in bullets. Native PowerPoint
placeholders/tables/pictures throughout (not one big image per slide), so
titles, bullets, tables, and pictures can each be edited independently.
No citations, no reproducibility slide — those live in the report only.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parent.parent
FIG = REPO / "reports/figures"
OUT = REPO / "reports/SLIDES.pptx"

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
GREEN = RGBColor(0x00, 0x83, 0x00)
MAGENTA = RGBColor(0xC2, 0x5B, 0x84)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
HEADER_FILL = RGBColor(0xEE, 0xF4, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def _set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SURFACE


def add_title_box(slide, text, top=Inches(0.4), size=32, color=INK, bold=True,
                   left=Inches(0.6), width=None, height=Inches(1.0)):
    width = width or (SLIDE_W - Inches(1.2))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def add_rule(slide, top, color=BLUE):
    line = slide.shapes.add_shape(1, Inches(0.6), top, SLIDE_W - Inches(1.2), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def add_bullets(slide, items, top=Inches(1.6), left=Inches(0.7),
                 width=None, height=None, size=22, line_spacing=1.3):
    width = width or (SLIDE_W - Inches(1.4))
    height = height or (SLIDE_H - top - Inches(0.4))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            bold_lead, rest = item
            r1 = p.add_run()
            r1.text = bold_lead
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = INK
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = INK_SECONDARY
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = INK
            r.font.name = "Calibri"
    return box


def add_table(slide, headers, rows, top=Inches(1.8), left=Inches(0.9),
              width=None, height=None, col_widths=None, font_size=18,
              header_color=BLUE):
    width = width or (SLIDE_W - Inches(1.8))
    height = height or Inches(0.55 * (len(rows) + 1))
    n_cols = len(headers)
    gshape = slide.shapes.add_table(len(rows) + 1, n_cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(font_size)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = INK
    return gshape


def add_picture_fit(slide, path, top, left=None, max_w=None, max_h=None):
    from PIL import Image
    max_w = max_w or (SLIDE_W - Inches(1.2))
    max_h = max_h or (SLIDE_H - top - Inches(0.3))
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w, h = max_w, max_w / aspect
    if h > max_h:
        h = max_h
        w = max_h * aspect
    left = left if left is not None else (SLIDE_W - w) / 2
    slide.shapes.add_picture(str(path), left, top, width=Emu(int(w)), height=Emu(int(h)))
    return top + h


def add_callout(slide, big_text, small_text, top=Inches(2.6), color=BLUE):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = big_text
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = small_text
    r2.font.size = Pt(20)
    r2.font.color.rgb = INK_SECONDARY
    r2.font.name = "Calibri"


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    return slide


# ---------------------------------------------------------------------------
# 1. Title
# ---------------------------------------------------------------------------
s = new_slide()
box = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Ancestry-Informative SNP Panels"
r.font.size = Pt(40)
r.font.bold = True
r.font.color.rgb = INK
r.font.name = "Calibri"
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "for Within-East-Asian Population Discrimination"
r2.font.size = Pt(24)
r2.font.color.rgb = BLUE
r2.font.name = "Calibri"
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.4))
tf2 = sub.text_frame
tf2.word_wrap = True
for i, line in enumerate(["Han  ·  Japanese  ·  Southeast Asian", "1000 Genomes Phase 3"]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.size = Pt(18)
    r.font.color.rgb = INK_SECONDARY
    r.font.name = "Calibri"
add_rule(s, Inches(4.2))

# ---------------------------------------------------------------------------
# 2. The Problem
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "The Problem")
add_rule(s, Inches(1.35))
add_bullets(s, [
    ("Goal: ", "smallest SNP panel that separates Han, Japanese, and Southeast Asian"),
    ("Hard part: ", "all three share the same continental ancestry — differences are small, diffuse"),
], top=Inches(1.9), size=26)

# ---------------------------------------------------------------------------
# 3. Why It Matters
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Why It Matters")
add_rule(s, Inches(1.35))
add_bullets(s, [
    ("Forensics — ", "ancestry finer than continental panels allow"),
    ("Pharmacogenomics — ", "drug response varies within continental groups"),
    ("Population genetics — ", "cheap, compact markers for fine structure"),
], top=Inches(1.9), size=24)

# ---------------------------------------------------------------------------
# 4. The Data
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "The Data")
add_rule(s, Inches(1.35))
add_table(s, ["Group", "1000G source", "n"], [
    ["Han", "CHB + CHS", "208"],
    ["JPT", "JPT", "104"],
    ["SEA", "KHV + CDX", "192"],
], top=Inches(2.0), font_size=20)
add_bullets(s, ["504 samples · 1000 Genomes Phase 3"], top=Inches(4.3), size=20)

# ---------------------------------------------------------------------------
# 5. Overall System Diagram
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Overall System", size=30)
add_picture_fit(s, FIG / "fig1_system_overview.png", top=Inches(1.2), max_h=Inches(5.8))

# ---------------------------------------------------------------------------
# 6. Panel Finder Detail
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, '"Panel Finder" — Detail', size=30)
add_picture_fit(s, FIG / "fig2_panelfinder_detail.png", top=Inches(1.2), max_h=Inches(5.8))

# ---------------------------------------------------------------------------
# 7. Candidate Pools
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Three Candidate Pools")
add_rule(s, Inches(1.35))
add_table(s, ["Pool", "Built from", "Size"], [
    ["stat", "χ², JSD, AFD", "1,005"],
    ["FST", "Hudson F_ST", "2,508"],
    ["fst_stat", "stat, cascaded onto FST", "965"],
], top=Inches(2.0), font_size=20)
add_bullets(s, ["No pool assumed superior — checked against blind selection"],
            top=Inches(4.4), size=20)

# ---------------------------------------------------------------------------
# 8. Two Honest Estimates
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Two Honest Estimates")
add_rule(s, Inches(1.35))
add_bullets(s, [
    ("Blind Selection — ", "picks config per fold, never sees test data. Zero bias."),
    ("Best Fixed Configuration — ", "names the single best of 54 candidates. Mild bias."),
    ("Both rebuilt inside every CV fold — ", "leak-free by design"),
], top=Inches(1.9), size=24)

# ---------------------------------------------------------------------------
# 9. Trust Hierarchy
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "How Much to Trust Each Piece", size=28)
add_picture_fit(s, FIG / "fig3_winner_consistency.png", top=Inches(1.2), max_h=Inches(3.7))
add_table(s, ["Component", "Blind agrees w/ grid"], [
    ["Reducer (EN)", "17 / 18"],
    ["Classifier", "11 / 18"],
    ["Pool", "4 / 18"],
], top=Inches(5.1), left=Inches(3.3), width=Inches(6.7), font_size=18)

# ---------------------------------------------------------------------------
# 10. The Pool Surprise
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "The Pool Surprise", size=30)
add_picture_fit(s, FIG / "fig5_pool_divergence.png", top=Inches(1.1), max_h=Inches(5.4))
add_bullets(s, [("Hindsight loves FST. Blind selection trusts fst_stat.", "")],
            top=Inches(6.6), size=18)

# ---------------------------------------------------------------------------
# 11. Validating the Approach
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Does the Search Work?", size=30)
add_picture_fit(s, FIG / "fig6_accuracy_vs_n.png", top=Inches(1.2), max_h=Inches(5.8))

# ---------------------------------------------------------------------------
# 12. Committed Panels
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Committed Panels: N = 35, 50, 70")
add_rule(s, Inches(1.35))
add_table(s, ["N", "Blind Selection", "Best Fixed Config"], [
    ["35", "86.91%", "87.90%"],
    ["50", "88.28%", "91.67%"],
    ["70", "92.86%", "94.25%"],
], top=Inches(2.1), font_size=22)

# ---------------------------------------------------------------------------
# 13. Parity with Cai-34
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Parity with Expert Curation", size=28)
add_callout(s, "94.25% vs 94.6%", "70-SNP automated panel  ·  vs. Cai-34 (34 SNPs, hand-curated)",
            top=Inches(2.8), color=BLUE)

# ---------------------------------------------------------------------------
# 14. Confusion Matrices
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Where Errors Happen", size=30)
add_picture_fit(s, FIG / "fig7_confusion_matrices.png", top=Inches(1.2), max_h=Inches(5.5))
add_bullets(s, [("JPT is hardest — confused with Han, not SEA", "")], top=Inches(6.9), size=18)

# ---------------------------------------------------------------------------
# 15. Recommendation
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Recommendation")
add_rule(s, Inches(1.35))
add_bullets(s, [
    ("ElasticNet — ", "always, no exception"),
    ("fst_stat — ", "default pool, trust blind selection over hindsight"),
    ("Pre-specify the classifier — ", "don't pick it post hoc"),
], top=Inches(1.9), size=26)

# ---------------------------------------------------------------------------
# 16. Conclusion
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Conclusion", size=32)
add_callout(s, "Parity, automated, transparent",
            "Every source of optimism measured and disclosed — not assumed away",
            top=Inches(3.0), color=GREEN)

# ---------------------------------------------------------------------------
# 17. Limitations
# ---------------------------------------------------------------------------
s = new_slide()
add_title_box(s, "Limitations")
add_rule(s, Inches(1.35))
add_bullets(s, [
    "Shipped panel's SNPs are chosen on all 504 samples",
    "54-candidate menu, not exhaustive",
    "Single cohort — generalizability untested",
], top=Inches(1.9), size=24)

prs.save(OUT)
print(f"Wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")
